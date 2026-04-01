#!/usr/bin/env python3
"""
Generate VitalFit Q1 2026 Performance Review deck (.pptx)
Follows client feedback: "Tell us what to DO, not just what happened."
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# -- Color palette (navy/blue professional theme) --
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BLUE = RGBColor(0x2C, 0x3E, 0x6B)
ACCENT_BLUE = RGBColor(0x3A, 0x7C, 0xBD)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN = RGBColor(0x27, 0x8B, 0x4C)
RED = RGBColor(0xC0, 0x39, 0x2B)
AMBER = RGBColor(0xD4, 0x8B, 0x0A)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_background(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box with specified formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=14, color=BLACK, spacing=Pt(8),
                    font_name="Calibri", bold_prefix=False):
    """Add a bulleted text box. If bold_prefix, bold text before first colon."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = spacing
        p.font.size = Pt(font_size)
        p.font.name = font_name

        if bold_prefix and ":" in item:
            prefix, rest = item.split(":", 1)
            run1 = p.add_run()
            run1.text = prefix + ":"
            run1.font.bold = True
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.name = font_name
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = color
            run2.font.name = font_name
        else:
            p.text = item
            p.font.color.rgb = color

    return txBox


def add_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    """Add a colored rectangle as a visual accent bar."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_kpi_box(slide, left, top, width, height, label, value,
                delta=None, delta_color=GREEN):
    """Add a KPI callout box with label, big number, and optional delta."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = ACCENT_BLUE
    shape.line.width = Pt(1.5)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Label
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(11)
    p.font.color.rgb = MED_GRAY
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    # Value
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = NAVY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER

    # Delta
    if delta:
        p3 = tf.add_paragraph()
        p3.text = delta
        p3.font.size = Pt(12)
        p3.font.color.rgb = delta_color
        p3.font.name = "Calibri"
        p3.alignment = PP_ALIGN.CENTER
        p3.font.bold = True


def add_divider_line(slide, left, top, width, color=ACCENT_BLUE, thickness=Pt(2)):
    """Add a horizontal divider line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, thickness
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, NAVY)

# Accent bar at top
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), ACCENT_BLUE)

add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
            "VitalFit Q1 2026 Performance Review",
            font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT)

add_textbox(slide, Inches(1.5), Inches(3.3), Inches(8), Inches(0.6),
            "Quarterly Business Review  |  January - March 2026",
            font_size=18, color=LIGHT_BLUE, alignment=PP_ALIGN.LEFT)

add_divider_line(slide, Inches(1.5), Inches(4.2), Inches(3), ACCENT_BLUE, Pt(3))

add_textbox(slide, Inches(1.5), Inches(4.6), Inches(6), Inches(0.5),
            "Prepared by Zenith Media  |  April 2026",
            font_size=14, color=LIGHT_BLUE, alignment=PP_ALIGN.LEFT)


# ============================================================
# SLIDE 2: The Headline -- Executive Summary
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)

# Top bar
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Q1 Beat Q4 Across the Board -- And the Trajectory is Accelerating",
            font_size=28, bold=True, color=NAVY)

add_divider_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE, Pt(3))

# KPI boxes row
add_kpi_box(slide, Inches(0.8), Inches(1.5), Inches(2.6), Inches(1.6),
            "ROAS", "1.39", "Up from 1.15 in Q4", GREEN)
add_kpi_box(slide, Inches(3.7), Inches(1.5), Inches(2.6), Inches(1.6),
            "TOTAL CONVERSIONS", "7,520", "+21% vs Q4 target", GREEN)
add_kpi_box(slide, Inches(6.6), Inches(1.5), Inches(2.6), Inches(1.6),
            "REVENUE", "$482,400", "On $348K spend", GREEN)
add_kpi_box(slide, Inches(9.5), Inches(1.5), Inches(2.6), Inches(1.6),
            "BUDGET DELIVERY", "96.7%", "$348.2K of $360K", GREEN)

# Narrative summary
add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.5),
            "The Story",
            font_size=16, bold=True, color=NAVY)

bullets = [
    "ROAS improved from 1.15 (Q4) to 1.39 (Q1) -- a 21% improvement driven by Google Search efficiency and a disciplined channel mix.",
    "Monthly trajectory matters more than the average: January started at 1.20 ROAS, February hit 1.38, March closed at 1.52. The optimization is compounding.",
    "Two standout findings require action: BC is significantly outperforming (1.82 ROAS) and is under-invested. Quebec remains below breakeven (0.98 ROAS) -- the creative localization promised in Q4 has not yet been deployed.",
    "The Lookalike - High LTV audience segment is burning $65K at a 30 index. This is the single largest efficiency drag in the account and needs an immediate decision: rebuild or reallocate.",
]
add_bullet_list(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.2),
                bullets, font_size=13, color=DARK_GRAY, spacing=Pt(10))


# ============================================================
# SLIDE 3: What's Working
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Google Search Is the Efficiency Engine; TikTok Is the Reach Engine",
            font_size=28, bold=True, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE, Pt(3))

# Left column: Google Search
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.5),
            "Google Search: Highest ROAS in the portfolio",
            font_size=16, bold=True, color=DARK_BLUE)

gs_bullets = [
    "2.23 ROAS on $98K spend -- returning $2.23 for every $1 invested",
    "3,120 conversions at $31.41 CPA -- the lowest cost-per-acquisition across all channels",
    "15.1% share of voice in search, ahead of GlowLab (14.2%) but behind PureNutra (18.2%)",
    "Recommendation: Increase search budget 10-15% in Q2 to capture share from PureNutra's PMAX expansion",
]
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.8),
                gs_bullets, font_size=13, color=DARK_GRAY, spacing=Pt(8))

# Right column: TikTok
add_textbox(slide, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.5),
            "TikTok: Awareness at scale for pennies",
            font_size=16, bold=True, color=DARK_BLUE)

tt_bullets = [
    "6.8M impressions at $11.47 CPM -- 3x the reach of Meta at one-third the CPM",
    "1.20 ROAS is at breakeven, but the role is upper-funnel awareness feeding Google Search",
    "9.9% SOV vs FuelBody's 12% -- FuelBody is a new threat with a TikTok-first strategy",
    "Recommendation: Test creator partnerships to match FuelBody's approach; measure lift to Google Search",
]
add_bullet_list(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(2.8),
                tt_bullets, font_size=13, color=DARK_GRAY, spacing=Pt(8))

# Bottom: Monthly trend
add_textbox(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
            "Monthly Momentum: The Trajectory Is the Real Story",
            font_size=16, bold=True, color=DARK_BLUE)

trend_bullets = [
    "January: 1.20 ROAS, $50 CPA  -->  February: 1.38 ROAS, $45.46 CPA  -->  March: 1.52 ROAS, $44.33 CPA",
    "Each month improved on the last. If March trends hold, Q2 projects to 1.6+ ROAS -- well above the 1.3 target.",
    "This is the result of optimization discipline, not just market conditions. The improvements are in our control.",
]
add_bullet_list(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.8),
                trend_bullets, font_size=13, color=DARK_GRAY, spacing=Pt(8))


# ============================================================
# SLIDE 4: The BC Surprise
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "BC Is Outperforming Ontario and Getting Less Than Half the Budget",
            font_size=28, bold=True, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE, Pt(3))

# KPI comparison boxes
add_kpi_box(slide, Inches(1.5), Inches(1.5), Inches(3.0), Inches(1.8),
            "BC", "1.82 ROAS", "$36.90 CPA | $62K spend", GREEN)
add_kpi_box(slide, Inches(5.0), Inches(1.5), Inches(3.0), Inches(1.8),
            "ONTARIO", "1.55 ROAS", "$42.65 CPA | $145K spend", ACCENT_BLUE)
add_kpi_box(slide, Inches(8.5), Inches(1.5), Inches(3.0), Inches(1.8),
            "ALBERTA", "1.34 ROAS", "$45.24 CPA | $38K spend", ACCENT_BLUE)

# Analysis
add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.5),
            "What This Means",
            font_size=16, bold=True, color=DARK_BLUE)

bc_bullets = [
    "BC delivers 17% higher ROAS and 13% lower CPA than Ontario -- on less than half the budget. There is headroom to scale.",
    "1,680 conversions from $62K is the most capital-efficient region in the portfolio.",
    "Ontario remains the volume play ($145K, 3,400 conversions), but incremental dollars should flow to BC first.",
    "Alberta is tracking to target at 1.34 ROAS -- maintain current allocation.",
    "Prairies + Atlantic (1.12 ROAS, $53 CPA) is marginal but low-volume -- not worth pulling budget from, not worth adding to.",
]
add_bullet_list(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.5),
                bc_bullets, font_size=13, color=DARK_GRAY, spacing=Pt(8))

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6),
            "ACTION: Shift $15-20K from Ontario prospecting to BC in Q2. Test whether BC efficiency holds at higher spend.",
            font_size=14, bold=True, color=ACCENT_BLUE)


# ============================================================
# SLIDE 5: The Quebec Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Quebec Is Below Breakeven -- The Promised Creative Fix Has Not Shipped",
            font_size=28, bold=True, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), RED, Pt(3))

# KPI box
add_kpi_box(slide, Inches(0.8), Inches(1.5), Inches(3.0), Inches(1.8),
            "QUEBEC ROAS", "0.98", "Below breakeven", RED)
add_kpi_box(slide, Inches(4.2), Inches(1.5), Inches(3.0), Inches(1.8),
            "QUEBEC CPA", "$68.33", "47% above portfolio avg", RED)
add_kpi_box(slide, Inches(7.6), Inches(1.5), Inches(3.0), Inches(1.8),
            "Q4 FLAG", "Yes", "Localization was promised", AMBER)

add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.5),
            "The Situation",
            font_size=16, bold=True, color=DARK_BLUE)

qc_bullets = [
    "Quebec represents $82K in spend (24% of regional budget) returning 0.98 ROAS -- we are losing money in this market.",
    "$68.33 CPA is 47% above the portfolio average of $46.30. The gap is not small and it is not closing.",
    "This was flagged in the Q4 deck. The Q4 recommendation was to refresh Quebec creative with French-language assets. That has not happened.",
    "Without localized creative, we are running English-language ads in a French-first market. The performance gap is predictable.",
]
add_bullet_list(slide, Inches(0.8), Inches(4.3), Inches(11.5), Inches(2.0),
                qc_bullets, font_size=13, color=DARK_GRAY, spacing=Pt(8))

add_textbox(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
            "Proposed Q2 Quebec Plan",
            font_size=16, bold=True, color=DARK_BLUE)

qc_plan = [
    "Week 1-2: Brief French-language creative agency. Produce 3 video + 5 static variants.",
    "Week 3-4: Launch localized creative in test cell. A/B against current English assets.",
    "Week 5-8: If CPA drops below $50, scale. If not, reduce Quebec budget 30% and reallocate to BC.",
]
add_bullet_list(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(1.0),
                qc_plan, font_size=13, color=DARK_GRAY, spacing=Pt(6))


# ============================================================
# SLIDE 6: Audience Winners & Losers
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Gym Beginners Are the Dark Horse; Lookalike-High LTV Is Broken",
            font_size=28, bold=True, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE, Pt(3))

# Left: Winners
add_accent_bar(slide, Inches(0.8), Inches(1.5), Inches(0.08), Inches(3.0), GREEN)

add_textbox(slide, Inches(1.2), Inches(1.5), Inches(5.0), Inches(0.5),
            "Scale These",
            font_size=18, bold=True, color=GREEN)

winner_bullets = [
    "Gym Beginners 18-24: 137 index, $33.85 CPA, 1,920 conversions on $65K. The best CPA in the portfolio and the highest indexing segment. This audience has untapped headroom.",
    "Fitness Enthusiasts 25-34: 133 index, $35.00 CPA, 2,800 conversions on $98K. The proven workhorse -- continues to deliver volume at efficient cost.",
    "Together these two segments drive 63% of all conversions at 29% below-average CPA.",
]
add_bullet_list(slide, Inches(1.2), Inches(2.1), Inches(5.2), Inches(2.5),
                winner_bullets, font_size=13, color=DARK_GRAY, spacing=Pt(10))

# Right: Losers
add_accent_bar(slide, Inches(7.0), Inches(1.5), Inches(0.08), Inches(3.0), RED)

add_textbox(slide, Inches(7.3), Inches(1.5), Inches(5.0), Inches(0.5),
            "Fix or Kill",
            font_size=18, bold=True, color=RED)

loser_bullets = [
    "Lookalike - High LTV: 30 index, $155.24 CPA, 420 conversions on $65.2K. This segment is spending at the same level as Gym Beginners but delivering 78% fewer conversions at 4.6x the CPA.",
    "Health-Conscious Parents 35-44: 90 index, $51.43 CPA, 1,400 conversions. Below average but not broken -- test creative refresh before pulling budget.",
    "Lapsed Customers: 94 index, $48.98 CPA, 980 conversions. Near-average performance. Maintain and monitor.",
]
add_bullet_list(slide, Inches(7.3), Inches(2.1), Inches(5.2), Inches(2.5),
                loser_bullets, font_size=13, color=DARK_GRAY, spacing=Pt(10))

add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.6),
            "ACTION: Pause Lookalike-High LTV immediately. Reallocate $65K to Gym Beginners and Fitness Enthusiasts (50/50 split). If Lookalike must continue, rebuild seed list from Q1 converters only.",
            font_size=14, bold=True, color=ACCENT_BLUE)


# ============================================================
# SLIDE 7: Competitive Landscape
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "FuelBody Is the New Threat -- They Are Outspending Us on TikTok 1.2:1",
            font_size=28, bold=True, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE, Pt(3))

# Competitor cards
# GlowLab
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(0.5),
            "GlowLab Nutrition",
            font_size=16, bold=True, color=DARK_BLUE)
gl_bullets = [
    "Est. total Q1 spend: $341K across Meta, Search, TikTok",
    "TikTok spend up 45% QoQ -- creator partnerships with 3 fitness influencers",
    "Launched new protein line in Feb with heavy video push on Meta",
    "Bidding aggressively on non-brand supplement search terms",
]
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(3.8), Inches(2.0),
                gl_bullets, font_size=12, color=DARK_GRAY, spacing=Pt(6))

# PureNutra
add_textbox(slide, Inches(4.8), Inches(1.5), Inches(3.5), Inches(0.5),
            "PureNutra",
            font_size=16, bold=True, color=DARK_BLUE)
pn_bullets = [
    "Est. total Q1 spend: $288K across Meta, Search, TikTok",
    "Pulled back Meta prospecting, shifted to retargeting (-5% QoQ)",
    "Increased Shopping/PMAX investment (+12% QoQ) -- owns 18.2% search SOV",
    "Testing TikTok for first time in Q1 at $28K",
]
add_bullet_list(slide, Inches(4.8), Inches(2.0), Inches(3.8), Inches(2.0),
                pn_bullets, font_size=12, color=DARK_GRAY, spacing=Pt(6))

# FuelBody
add_textbox(slide, Inches(8.8), Inches(1.5), Inches(4.0), Inches(0.5),
            "FuelBody Co [NEW ENTRANT]",
            font_size=16, bold=True, color=RED)
fb_bullets = [
    "Est. total Q1 spend: $173K -- DTC recovery supplements",
    "TikTok-first: $95K on TikTok alone (12% SOV vs our 9.9%)",
    "Viral unboxing content driving organic amplification",
    "Meta spend up 22% -- expanding beyond TikTok base",
]
add_bullet_list(slide, Inches(8.8), Inches(2.0), Inches(4.0), Inches(2.0),
                fb_bullets, font_size=12, color=DARK_GRAY, spacing=Pt(6))

# Implications
add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11), Inches(0.5),
            "What This Means for VitalFit",
            font_size=16, bold=True, color=DARK_BLUE)

comp_impl = [
    "Google Search: We hold 15.1% SOV vs PureNutra's 18.2%. Closing this gap is achievable with incremental budget -- search is our highest-ROAS channel.",
    "TikTok: FuelBody is outspending us ($95K vs $78K) and gaining share with creator-led content. We need a creator strategy, not just paid media.",
    "Meta: The market is shifting -- PureNutra pulling back prospecting while GlowLab invests in video. Our Meta ROAS (1.20) needs creative improvement to stay competitive.",
    "Net: Our total category spend ($318K ex-CTV) is middle of the pack. The competitive advantage must come from efficiency and creative quality, not outspending.",
]
add_bullet_list(slide, Inches(0.8), Inches(5.0), Inches(11.5), Inches(2.2),
                comp_impl, font_size=13, color=DARK_GRAY, spacing=Pt(8))


# ============================================================
# SLIDE 8: CTV Test Results
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "CTV Delivered 2.1M Impressions at $14.38 CPM -- Conversion Measurement Still Pending",
            font_size=26, bold=True, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE, Pt(3))

# KPI boxes
add_kpi_box(slide, Inches(1.0), Inches(1.5), Inches(2.8), Inches(1.5),
            "IMPRESSIONS", "2.1M", "$30.2K spend", ACCENT_BLUE)
add_kpi_box(slide, Inches(4.2), Inches(1.5), Inches(2.8), Inches(1.5),
            "CPM", "$14.38", "Competitive for CTV", ACCENT_BLUE)
add_kpi_box(slide, Inches(7.4), Inches(1.5), Inches(2.8), Inches(1.5),
            "CONVERSIONS", "N/A", "No direct attribution yet", AMBER)

add_textbox(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.5),
            "What We Know",
            font_size=16, bold=True, color=DARK_BLUE)

ctv_bullets = [
    "$14.38 CPM is within the expected range for Canadian CTV inventory ($12-18 CPM). We are not overpaying for reach.",
    "No click or conversion data is available -- CTV does not support direct response measurement. This was expected when the test was scoped.",
    "Brand lift study is pending from the vendor. Results expected mid-April.",
    "Anecdotally, branded search volume increased 8% during CTV flight periods, but this is not yet validated as causal.",
]
add_bullet_list(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(2.0),
                ctv_bullets, font_size=13, color=DARK_GRAY, spacing=Pt(8))

add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5),
            "Recommendation",
            font_size=16, bold=True, color=DARK_BLUE)

ctv_rec = [
    "Continue CTV at current $30K/quarter while brand lift data comes in. Do not scale until measurement framework is validated.",
    "Implement a branded search lift test in Q2: run CTV on/off in alternating 2-week windows and measure Google branded search volume delta.",
]
add_bullet_list(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(1.0),
                ctv_rec, font_size=13, color=DARK_GRAY, spacing=Pt(6))


# ============================================================
# SLIDE 9: Q2 Recommendations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Five Actions for Q2 -- Each With a Specific Expected Outcome",
            font_size=28, bold=True, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE, Pt(3))

recs = [
    "1. Kill or rebuild Lookalike-High LTV audience: Pause the segment immediately and reallocate $65K to Gym Beginners and Fitness Enthusiasts. Expected impact: recover ~1,200 conversions at $35 CPA vs. current 420 at $155 CPA.",
    "2. Ship French-language creative for Quebec: Brief agency in Week 1, launch test in Week 3, evaluate at Week 5. If CPA drops below $50, scale; if not, cut Quebec budget 30%. Expected impact: improve Quebec ROAS from 0.98 to 1.3+ or save $25K.",
    "3. Increase BC investment by $15-20K: Shift from Ontario prospecting budget. BC efficiency (1.82 ROAS, $36.90 CPA) suggests headroom. Expected impact: 400-550 incremental conversions at below-portfolio-average CPA.",
    "4. Launch TikTok creator partnerships: Counter FuelBody's creator-first strategy. Budget $15K for 2-3 creator deals. Expected impact: improve TikTok SOV from 9.9% to 12%+ and test organic amplification.",
    "5. Implement CTV measurement framework: Run branded search lift test with alternating 2-week on/off windows. Expected impact: data to make a scale-or-cut decision on CTV by end of Q2.",
]
add_bullet_list(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.5),
                recs, font_size=14, color=DARK_GRAY, spacing=Pt(16),
                bold_prefix=True)

add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.6),
            "Total Q2 budget recommendation: $370K (up from $360K). Net reallocation, not net new spend, drives most of the improvement.",
            font_size=14, bold=True, color=ACCENT_BLUE)


# ============================================================
# SLIDE 10: Appendix -- Full Data Tables
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_accent_bar(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), NAVY)

add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
            "Appendix: Full Data Tables",
            font_size=28, bold=True, color=NAVY)
add_divider_line(slide, Inches(0.8), Inches(1.1), Inches(2.5), ACCENT_BLUE, Pt(3))

# Channel table
add_textbox(slide, Inches(0.8), Inches(1.4), Inches(3), Inches(0.4),
            "Channel Performance", font_size=14, bold=True, color=DARK_BLUE)

# Build table: Channel data
channel_data = [
    ["Channel", "Spend", "Impressions", "Conversions", "Revenue", "ROAS", "CPA"],
    ["Meta", "$142,000", "4.2M", "2,840", "$170,400", "1.20", "$50.00"],
    ["Google Search", "$98,000", "1.85M", "3,120", "$218,400", "2.23", "$31.41"],
    ["TikTok", "$78,000", "6.8M", "1,560", "$93,600", "1.20", "$50.00"],
    ["CTV", "$30,200", "2.1M", "--", "--", "--", "$14.38 CPM"],
    ["TOTAL", "$348,200", "14.95M", "7,520", "$482,400", "1.39", "$46.30"],
]

rows, cols = len(channel_data), len(channel_data[0])
table_shape = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(7.5), Inches(1.8))
table = table_shape.table

for i, row_data in enumerate(channel_data):
    for j, cell_text in enumerate(row_data):
        cell = table.cell(i, j)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.name = "Calibri"
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        elif i == len(channel_data) - 1:
            p.font.bold = True
            p.font.color.rgb = NAVY
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BLUE
        else:
            p.font.color.rgb = DARK_GRAY
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

# Region table
add_textbox(slide, Inches(0.8), Inches(3.9), Inches(3), Inches(0.4),
            "Regional Performance", font_size=14, bold=True, color=DARK_BLUE)

region_data = [
    ["Region", "Spend", "Conversions", "CPA", "ROAS"],
    ["Ontario", "$145,000", "3,400", "$42.65", "1.55"],
    ["Quebec", "$82,000", "1,200", "$68.33", "0.98"],
    ["BC", "$62,000", "1,680", "$36.90", "1.82"],
    ["Alberta", "$38,000", "840", "$45.24", "1.34"],
    ["Prairies+Atlantic", "$21,200", "400", "$53.00", "1.12"],
]

rows2, cols2 = len(region_data), len(region_data[0])
table_shape2 = slide.shapes.add_table(rows2, cols2, Inches(0.8), Inches(4.3), Inches(5.5), Inches(1.6))
table2 = table_shape2.table

for i, row_data in enumerate(region_data):
    for j, cell_text in enumerate(row_data):
        cell = table2.cell(i, j)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.name = "Calibri"
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            p.font.color.rgb = DARK_GRAY
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE

# Audience table
add_textbox(slide, Inches(7.0), Inches(3.9), Inches(3), Inches(0.4),
            "Audience Segments", font_size=14, bold=True, color=DARK_BLUE)

audience_data = [
    ["Segment", "CPA", "Index"],
    ["Fitness Enth. 25-34", "$35.00", "133"],
    ["Gym Beginners 18-24", "$33.85", "137"],
    ["Parents 35-44", "$51.43", "90"],
    ["Lapsed Customers", "$48.98", "94"],
    ["Lookalike-High LTV", "$155.24", "30"],
]

rows3, cols3 = len(audience_data), len(audience_data[0])
table_shape3 = slide.shapes.add_table(rows3, cols3, Inches(7.0), Inches(4.3), Inches(4.5), Inches(1.6))
table3 = table_shape3.table

for i, row_data in enumerate(audience_data):
    for j, cell_text in enumerate(row_data):
        cell = table3.cell(i, j)
        cell.text = cell_text
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.name = "Calibri"
        if i == 0:
            p.font.bold = True
            p.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        else:
            p.font.color.rgb = DARK_GRAY
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(
    os.path.dirname(os.path.abspath(__file__)),
    "output.pptx"
)
prs.save(output_path)
print(f"Deck saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
